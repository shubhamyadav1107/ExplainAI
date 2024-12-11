# app.py

import streamlit as st
import os
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from transformers import pipeline
from graphviz import Digraph
from PIL import Image

# Initialize the summarization pipeline
@st.cache_resource
def load_summarizer():
    return pipeline("summarization", model="facebook/bart-large-cnn")

summarizer = load_summarizer()

# Function to extract text from PDF
def extract_text_from_pdf(file_path):
    reader = PdfReader(file_path)
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
    return text

# Function to extract text from Word document
def extract_text_from_word(file_path):
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

# Function to extract text from PPT
def extract_text_from_ppt(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text + "\n"
    return text

# General function to handle extraction
def extract_text(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext == ".docx":
        return extract_text_from_word(file_path)
    elif ext == ".pptx":
        return extract_text_from_ppt(file_path)
    else:
        return "Unsupported file format."

# Function for summarization
def summarize_text(text, chunk_size=1000, max_length=150, min_length=50):
    chunks = [text[i:i + chunk_size] for i in range(0, len(text), chunk_size)]
    summary = ""
    for chunk in chunks:
        summarized = summarizer(chunk, max_length=max_length, min_length=min_length, do_sample=False)
        summary += summarized[0]['summary_text'] + " "
    return summary.strip()

# Function to create a flowchart
def create_flowchart(summarized_text, output_file="flowchart"):
    points = summarized_text.split('. ')
    dot = Digraph(comment="Visual Representation", format='png')
    for i, point in enumerate(points):
        if point.strip():
            dot.node(f'Node{i}', point.strip(), shape='box', style='filled', color='lightblue')
            if i > 0:
                dot.edge(f'Node{i-1}', f'Node{i}')
    dot.render(output_file, cleanup=True)
    return f"{output_file}.png"

# Streamlit UI
st.title("📝🤓ExplainAI👽: Transforming Learning With AI")

# File upload
uploaded_file = st.file_uploader("Upload a file (PDF, DOCX, PPTX)", type=["pdf", "docx", "pptx"])

if uploaded_file:
    # Save the uploaded file temporarily
    file_path = os.path.join("temp", uploaded_file.name)
    os.makedirs("temp", exist_ok=True)
    with open(file_path, "wb") as f:
        f.write(uploaded_file.getvalue())

    # Extract text
    st.subheader("🔍 Extracting Text...")
    extracted_text = extract_text(file_path)
    st.text_area("Extracted Text:", extracted_text, height=300)

    # Summarize text
    st.subheader("✍️ Summarizing Text...")
    summarized_text = summarize_text(extracted_text)
    st.text_area("Summarized Text:", summarized_text, height=200)

    # Generate flowchart
    st.subheader("🖼️ Generating Flowchart...")
    flowchart_path = create_flowchart(summarized_text, output_file="flowchart")
    st.image(flowchart_path, caption="Generated Flowchart", use_column_width=True)

    # Clean up
    os.remove(file_path)
